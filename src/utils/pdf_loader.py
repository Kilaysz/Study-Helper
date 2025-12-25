import io
import os
import pymupdf4llm
from pptx import Presentation

# --- IN-MEMORY CACHE ---
# Even though we save files to disk, we keep a cache for the 
# current session's text so nodes don't have to re-read files constantly.
_pdf_cache = {}

def save_to_cache(thread_id: str, text: str):
    """Stores text in the global in-memory cache for fast access by nodes."""
    if thread_id:
        print(f"🧠 Cache: Updating context for thread {thread_id} with newest document.")
        _pdf_cache[thread_id] = text

def get_from_cache(thread_id: str) -> str:
    """Retrieves text from the global in-memory cache."""
    return _pdf_cache.get(thread_id, "")

def clear_cache(thread_id: str):
    """Removes a specific entry from the cache."""
    if thread_id in _pdf_cache:
        del _pdf_cache[thread_id]

# --- LOADING LOGIC ---

def load_pptx_content(file_path: str) -> str:
    """
    Extracts text from a PowerPoint file saved on disk.
    """
    try:
        ppt = Presentation(file_path)
        text_runs = []
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e:
        print(f"Error extracting PPTX from {file_path}: {e}")
        return ""

def load_pdf_content(file_path: str) -> str:
    """
    Primary function for loading documents from the 'data/' folder.
    Determines parser (PDF or PPTX) based on file extension.
    """
    if not os.path.exists(file_path):
        print(f"Error: File not found at {file_path}")
        return ""

    try:
        # Check extension for PPTX
        if file_path.lower().endswith(".pptx"):
            return load_pptx_content(file_path)
        
        # Default to PDF (optimized for Markdown extraction)
        return pymupdf4llm.to_markdown(file_path)
    except Exception as e:
        print(f"Error loading document from path {file_path}: {e}")
        return ""

def load_pdf_content_from_bytes(file_bytes: bytes, file_type: str = "pdf") -> str:
    """
    Fallback utility to convert document bytes into text.
    """
    try:
        if file_type.lower() == "pptx":
            ppt = Presentation(io.BytesIO(file_bytes))
            text_runs = []
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
        
        # Default to PDF
        pdf_stream = io.BytesIO(file_bytes)
        return pymupdf4llm.to_markdown(pdf_stream)
    except Exception as e:
        print(f"Error extracting document text from memory: {e}")
        return ""